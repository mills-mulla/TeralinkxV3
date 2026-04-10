# apps/finance/board_report_export.py
"""
Board Report Export Service
Handles PDF and PowerPoint export for board reports.
"""
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_RIGHT
from io import BytesIO
from django.core.mail import EmailMessage
from django.conf import settings
import logging

logger = logging.getLogger(__name__)


class BoardReportExporter:
    """Export board reports to PDF and PowerPoint formats"""
    
    @staticmethod
    def export_to_pdf(report):
        """Generate PDF export of board report"""
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=0.5*inch, bottomMargin=0.5*inch)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a1a'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2c3e50'),
            spaceAfter=12,
            spaceBefore=20
        )
        
        # Title
        story.append(Paragraph(f"Board Report - {report.report_period_display}", title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # Executive Summary
        story.append(Paragraph("Executive Summary", heading_style))
        story.append(Paragraph(report.executive_summary.replace('\n', '<br/>'), styles['Normal']))
        story.append(Spacer(1, 0.2*inch))
        
        # Financial Performance
        story.append(Paragraph("Financial Performance", heading_style))
        fp = report.financial_performance
        
        financial_data = [
            ['Metric', 'Current', 'Previous', 'Change'],
            ['Revenue', f"KES {fp['revenue']['current']:,.2f}", 
             f"KES {fp['revenue']['previous']:,.2f}", 
             f"{fp['revenue']['growth_pct']:+.1f}%"],
            ['Expenses', f"KES {fp['expenses']['current']:,.2f}", 
             f"KES {fp['expenses']['previous']:,.2f}", 
             f"{fp['expenses']['change_pct']:+.1f}%"],
            ['Net Profit', f"KES {fp['net_profit']:,.2f}", '', 
             f"{fp['profit_margin_pct']:.1f}% margin"],
        ]
        
        financial_table = Table(financial_data, colWidths=[2*inch, 1.5*inch, 1.5*inch, 1.5*inch])
        financial_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498db')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(financial_table)
        story.append(Spacer(1, 0.2*inch))
        
        # Customer Metrics
        story.append(Paragraph("Customer Metrics", heading_style))
        cm = report.customer_metrics
        
        customer_data = [
            ['Metric', 'Value'],
            ['Active Customers', f"{cm['active_customers']:,}"],
            ['New Customers', f"{cm['new_customers']:,}"],
            ['Churned Customers', f"{cm['churned_customers']:,}"],
            ['Churn Rate', f"{cm['churn_rate_pct']:.1f}%"],
            ['ARPU', f"KES {cm['arpu']:,.2f}"],
        ]
        
        customer_table = Table(customer_data, colWidths=[3*inch, 3*inch])
        customer_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2ecc71')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(customer_table)
        story.append(Spacer(1, 0.2*inch))
        
        # Key Highlights
        if report.key_highlights:
            story.append(Paragraph("Key Highlights", heading_style))
            for highlight in report.key_highlights:
                story.append(Paragraph(f"• {highlight}", styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        # Challenges
        if report.challenges:
            story.append(Paragraph("Challenges", heading_style))
            for challenge in report.challenges:
                story.append(Paragraph(f"• {challenge}", styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        # Recommendations
        if report.recommendations:
            story.append(Paragraph("Recommendations", heading_style))
            for recommendation in report.recommendations:
                story.append(Paragraph(f"• {recommendation}", styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        # Risk Register
        story.append(PageBreak())
        story.append(Paragraph("Risk Register", heading_style))
        rr = report.risk_register
        
        risk_data = [
            ['Risk Category', 'Value'],
            ['High Risk Customers', f"{rr['high_risk_customers']:,}"],
            ['Revenue at Risk', f"KES {rr['revenue_at_risk']:,.2f}"],
            ['Outstanding Receivables', f"KES {rr['outstanding_receivables']:,.2f}"],
        ]
        
        risk_table = Table(risk_data, colWidths=[3*inch, 3*inch])
        risk_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#e74c3c')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.lightgrey),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        story.append(risk_table)
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def export_to_pptx(report):
        """Generate PowerPoint export of board report"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            raise ImportError("python-pptx library required for PowerPoint export")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = f"Board Report"
        subtitle.text = f"{report.report_period_display}\nGenerated: {report.created_at.strftime('%B %d, %Y')}"
        
        # Executive Summary Slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Executive Summary"
        tf = body_shape.text_frame
        tf.text = report.executive_summary
        
        # Financial Performance Slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Financial Performance"
        
        fp = report.financial_performance
        
        # Add table
        rows, cols = 4, 4
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header row
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Current"
        table.cell(0, 2).text = "Previous"
        table.cell(0, 3).text = "Change"
        
        # Data rows
        table.cell(1, 0).text = "Revenue"
        table.cell(1, 1).text = f"KES {fp['revenue']['current']:,.2f}"
        table.cell(1, 2).text = f"KES {fp['revenue']['previous']:,.2f}"
        table.cell(1, 3).text = f"{fp['revenue']['growth_pct']:+.1f}%"
        
        table.cell(2, 0).text = "Expenses"
        table.cell(2, 1).text = f"KES {fp['expenses']['current']:,.2f}"
        table.cell(2, 2).text = f"KES {fp['expenses']['previous']:,.2f}"
        table.cell(2, 3).text = f"{fp['expenses']['change_pct']:+.1f}%"
        
        table.cell(3, 0).text = "Net Profit"
        table.cell(3, 1).text = f"KES {fp['net_profit']:,.2f}"
        table.cell(3, 2).text = ""
        table.cell(3, 3).text = f"{fp['profit_margin_pct']:.1f}% margin"
        
        # Customer Metrics Slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Customer Metrics"
        
        cm = report.customer_metrics
        
        rows, cols = 6, 2
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        
        table.cell(1, 0).text = "Active Customers"
        table.cell(1, 1).text = f"{cm['active_customers']:,}"
        
        table.cell(2, 0).text = "New Customers"
        table.cell(2, 1).text = f"{cm['new_customers']:,}"
        
        table.cell(3, 0).text = "Churned Customers"
        table.cell(3, 1).text = f"{cm['churned_customers']:,}"
        
        table.cell(4, 0).text = "Churn Rate"
        table.cell(4, 1).text = f"{cm['churn_rate_pct']:.1f}%"
        
        table.cell(5, 0).text = "ARPU"
        table.cell(5, 1).text = f"KES {cm['arpu']:,.2f}"
        
        # Key Highlights Slide
        if report.key_highlights:
            slide = prs.slides.add_slide(bullet_slide_layout)
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            
            title_shape.text = "Key Highlights"
            tf = body_shape.text_frame
            
            for highlight in report.key_highlights:
                p = tf.add_paragraph()
                p.text = highlight
                p.level = 0
        
        # Challenges & Recommendations Slide
        if report.challenges or report.recommendations:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = "Challenges & Recommendations"
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(4.5)
            height = Inches(5)
            
            # Challenges text box
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Challenges:"
            p.font.bold = True
            p.font.size = Pt(14)
            
            for challenge in report.challenges:
                p = tf.add_paragraph()
                p.text = f"• {challenge}"
                p.font.size = Pt(11)
            
            # Recommendations text box
            left = Inches(5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Recommendations:"
            p.font.bold = True
            p.font.size = Pt(14)
            
            for rec in report.recommendations:
                p = tf.add_paragraph()
                p.text = f"• {rec}"
                p.font.size = Pt(11)
        
        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    @staticmethod
    def send_board_report_email(report, recipient_emails, include_pdf=True, include_pptx=True):
        """Send board report via email with attachments"""
        subject = f"Board Report - {report.report_period_display}"
        
        body = f"""
Dear Board Members,

Please find attached the board report for {report.report_period_display}.

Executive Summary:
{report.executive_summary}

Key Highlights:
{chr(10).join(f'• {h}' for h in report.key_highlights)}

This report was automatically generated on {report.created_at.strftime('%B %d, %Y at %I:%M %p')}.

Best regards,
TeralinkX Finance Team
        """
        
        email = EmailMessage(
            subject=subject,
            body=body,
            from_email=settings.DEFAULT_FROM_EMAIL,
            to=recipient_emails
        )
        
        # Attach PDF
        if include_pdf:
            try:
                pdf_buffer = BoardReportExporter.export_to_pdf(report)
                email.attach(
                    f"board_report_{report.report_period_display}.pdf",
                    pdf_buffer.getvalue(),
                    'application/pdf'
                )
                logger.info(f"PDF attached to email for report {report.id}")
            except Exception as e:
                logger.error(f"Failed to attach PDF: {e}")
        
        # Attach PowerPoint
        if include_pptx:
            try:
                pptx_buffer = BoardReportExporter.export_to_pptx(report)
                email.attach(
                    f"board_report_{report.report_period_display}.pptx",
                    pptx_buffer.getvalue(),
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                )
                logger.info(f"PowerPoint attached to email for report {report.id}")
            except Exception as e:
                logger.error(f"Failed to attach PowerPoint: {e}")
        
        # Send email
        try:
            email.send()
            logger.info(f"Board report email sent to {len(recipient_emails)} recipients")
            return True
        except Exception as e:
            logger.error(f"Failed to send board report email: {e}")
            return False
